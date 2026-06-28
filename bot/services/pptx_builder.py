"""Build a .pptx file from the LLM's JSON slide description.

Input shape (legacy template path; the rich path is planned by DECK_PLAN_SYSTEM):
    {"title": "...", "slides": [{"title": "...", "bullets": [...], "image_ref": 1}]}

`image_ref` (optional) is a photo id from the batch; the matching image bytes are
placed on the slide. We tolerate the model wrapping the JSON in ```json fences or
adding stray text, so we extract the first {...} block before parsing.
"""

from __future__ import annotations

import io
import json
import logging
import re

from pptx import Presentation
from pptx.util import Pt

logger = logging.getLogger(__name__)

# Layout indexes in the default python-pptx template.
_TITLE_SLIDE = 0
_TITLE_AND_CONTENT = 1

_JSON_BLOCK = re.compile(r"\{.*\}", re.DOTALL)


def _place_image(slide, img_bytes: bytes, left: int, top: int, box_w: int, box_h: int,
                 center_in_box: bool = False) -> bool:
    """Add a picture sized to fit within (box_w, box_h), aspect preserved.

    We set width only, then scale both dimensions by the same ratio if it's too
    tall — never width and height independently, so the image never distorts.
    Returns True on success.
    """
    try:
        pic = slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=box_w)
    except Exception:  # noqa: BLE001 - bad/unsupported image bytes -> skip
        logger.warning("Skipped a slide image (could not be placed)")
        return False
    if pic.height > box_h:
        ratio = box_h / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = box_h
    if center_in_box:
        pic.left = left + (box_w - pic.width) // 2
        pic.top = top + (box_h - pic.height) // 2
    return True


def parse_slides(raw: str) -> dict:
    """Parse the model output into a {..., slides} dict (lenient).

    Works for both the legacy {title, slides:[{title,bullets}]} and the rich
    {palette, slides:[{layout,...}]} plan — both have a 'slides' array.
    """
    match = _JSON_BLOCK.search(raw or "")
    if not match:
        raise ValueError("no JSON object found in model output")
    data = json.loads(match.group(0))
    if not isinstance(data, dict) or "slides" not in data:
        raise ValueError("JSON missing 'slides'")
    return data


# --- Design-system orchestration (rich plan) ---------------------------------
def build_design(plan: dict, photos: dict | None = None, overrides: dict | None = None) -> bytes:
    """Render a rich layout plan with the design system (deck_design)."""
    from bot.services import deck_design  # local import keeps this module light

    return deck_design.build_deck(plan, photos=photos, overrides=overrides)


def flatten_for_template(plan: dict) -> dict:
    """Collapse a rich layout plan into the legacy {title, slides} shape so a
    company template (placeholder-based) can render it while keeping its theme."""
    deck_title = ""
    slides: list[dict] = []
    for s in plan.get("slides", []):
        if not isinstance(s, dict):
            continue
        layout = s.get("layout")
        title = str(s.get("title") or "")
        if layout == "title" and not deck_title:
            deck_title = title or str(s.get("subtitle") or "")
            continue
        bullets: list[str] = []
        if layout == "quote":
            bullets = [f"“{s.get('quote', '')}”", f"— {s.get('attribution', '')}"]
            title = title or "Quote"
        else:
            bullets += [str(b) for b in (s.get("bullets") or [])]
            bullets += [str(b) for b in (s.get("items") or [])]
            for st in (s.get("stats") or []):
                bullets.append(f"{st.get('value', '')} — {st.get('label', '')}")
            for col in (s.get("columns") or []):
                bullets.append(str(col.get("heading", "")))
                bullets += [f"  • {it}" for it in col.get("items", [])]
            if s.get("caption"):
                bullets.append(str(s["caption"]))
        out: dict = {"title": title, "bullets": [b for b in bullets if b.strip()]}
        if isinstance(s.get("image_ref"), int):
            out["image_ref"] = s["image_ref"]
        slides.append(out)
    return {"title": deck_title or "Presentation", "slides": slides}


def build_pptx(
    data: dict,
    template_bytes: bytes | None = None,
    photos: dict[int, bytes] | None = None,
    gallery: bool = False,
    gallery_title: str = "Images",
) -> bytes:
    """Render the slide dict into .pptx bytes.

    `photos` maps a photo id -> image bytes; a slide's `image_ref` places that
    photo. If `gallery` is True, every photo not already placed is appended as a
    one-per-slide gallery at the end, so all forwarded photos appear at least once.

    A `template_bytes` (.pptx/.potx) is used as the base so the company theme and
    layouts are preserved; falls back to the default theme if it can't be opened.
    """
    photos = photos or {}
    if template_bytes:
        try:
            prs = Presentation(io.BytesIO(template_bytes))
        except Exception:  # noqa: BLE001 - bad/unsupported template -> default
            prs = Presentation()
    else:
        prs = Presentation()

    slide_w, slide_h = prs.slide_width, prs.slide_height
    used_ids: set[int] = set()

    # Title slide.
    title_slide = prs.slides.add_slide(prs.slide_layouts[_TITLE_SLIDE])
    title_slide.shapes.title.text = str(data.get("title") or "Presentation")

    # Content slides.
    for slide in data.get("slides", []):
        s = prs.slides.add_slide(prs.slide_layouts[_TITLE_AND_CONTENT])
        s.shapes.title.text = str(slide.get("title") or "")

        bullets = [str(b) for b in slide.get("bullets", []) if str(b).strip()]
        body = s.placeholders[1].text_frame
        body.clear()
        for i, bullet in enumerate(bullets):
            paragraph = body.paragraphs[0] if i == 0 else body.add_paragraph()
            paragraph.text = bullet
            paragraph.font.size = Pt(24)
            paragraph.font.name = "Calibri"

        # Optional image for this slide.
        ref = slide.get("image_ref")
        img = photos.get(ref) if isinstance(ref, int) else None
        if img:
            if bullets:
                # Right ~40% of the slide so it sits beside the text.
                placed = _place_image(
                    s, img, int(slide_w * 0.56), int(slide_h * 0.26),
                    int(slide_w * 0.40), int(slide_h * 0.58), center_in_box=True,
                )
            else:
                # Image-only slide: larger and centered.
                placed = _place_image(
                    s, img, int(slide_w * 0.15), int(slide_h * 0.22),
                    int(slide_w * 0.70), int(slide_h * 0.66), center_in_box=True,
                )
            if placed:
                used_ids.add(ref)

    # Gallery: ensure every forwarded photo appears at least once.
    if gallery:
        for pid in sorted(photos):
            if pid in used_ids:
                continue
            _add_gallery_slide(prs, photos[pid], gallery_title, slide_w, slide_h)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _add_gallery_slide(prs, img_bytes: bytes, title: str, slide_w: int, slide_h: int) -> None:
    """A title + one centered photo (used for the fallback gallery)."""
    s = prs.slides.add_slide(prs.slide_layouts[_TITLE_AND_CONTENT])
    s.shapes.title.text = title
    # Drop the empty body placeholder so it doesn't sit under the image.
    for ph in list(s.placeholders):
        if ph.placeholder_format.idx == 1:
            ph._element.getparent().remove(ph._element)
    _place_image(
        s, img_bytes, int(slide_w * 0.15), int(slide_h * 0.22),
        int(slide_w * 0.70), int(slide_h * 0.66), center_in_box=True,
    )
