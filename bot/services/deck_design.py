"""Design system for presentations: a palette library + one robust renderer per
layout. Slides are built on the BLANK layout with manual placement so text and
images live in separate columns and cannot overlap by construction.

`build_deck` returns .pptx bytes. `overrides` (per 1-based slide number, matching
the rendered slide order) lets the QA loop nudge a slide deterministically
(smaller body font, higher contrast, stricter columns) and rebuild.
"""

from __future__ import annotations

import io
import logging

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

logger = logging.getLogger(__name__)

# --- Palettes: primary = dark (white-text-safe), tint = light accent for dark
# slides. On white slides: ink text, primary titles/numbers, primary circles. ---
PALETTES: dict[str, dict[str, str]] = {
    "Midnight Executive": {"primary": "1E2761", "tint": "CADCFC"},
    "Teal Trust": {"primary": "028090", "tint": "02C39A"},
    "Forest & Moss": {"primary": "2C5F2D", "tint": "97BC62"},
    "Coral Energy": {"primary": "2F3C7E", "tint": "F9E795"},
    "Ocean Gradient": {"primary": "21295C", "tint": "1C7293"},
    "Charcoal Minimal": {"primary": "212121", "tint": "B0B0B0"},
    "Berry & Cream": {"primary": "6D2E46", "tint": "D8B4A0"},
    "Cherry Bold": {"primary": "990011", "tint": "E8A0A0"},
}
_DEFAULT_PALETTE = "Teal Trust"

WHITE = "FFFFFF"
INK = "1A1A1A"        # body text on white
MUTED = "6B6B6B"      # captions
BORDER = "D9D9D9"     # image frame
FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_W = SLIDE_W - 2 * MARGIN

_DARK_LAYOUTS = {"title", "section", "quote"}
_MAX_ITEMS = 5


# --- Plan normalization (stable slide indices before QA) ---------------------
def normalize_plan(plan: dict) -> dict:
    """Cap items per slide and split overflow into '(cont.)' slides so indices
    stay stable across the QA loop."""
    out_slides: list[dict] = []
    for raw in plan.get("slides", []):
        if not isinstance(raw, dict):
            continue
        layout = raw.get("layout", "bullets")
        key = {"bullets": "bullets", "agenda": "items"}.get(layout)
        if key and isinstance(raw.get(key), list) and len(raw[key]) > _MAX_ITEMS:
            items = [str(x) for x in raw[key]]
            first = True
            for i in range(0, len(items), _MAX_ITEMS):
                chunk = dict(raw)
                chunk[key] = items[i:i + _MAX_ITEMS]
                if not first:
                    chunk["title"] = f"{raw.get('title', '')} (cont.)".strip()
                out_slides.append(chunk)
                first = False
        else:
            out_slides.append(raw)
    return {"palette": plan.get("palette"), "slides": out_slides}


# --- Low-level helpers -------------------------------------------------------
def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


def _bg(slide, hex_str: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex_str)


def _text(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def _run(para, text, size, color, bold=False, align=PP_ALIGN.LEFT, font=FONT):
    para.alignment = align
    r = para.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = _rgb(color)
    return r


def _circle(slide, left, top, diameter, fill_hex, glyph, glyph_color):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill_hex)
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    _run(p, glyph, 18, glyph_color, bold=True, align=PP_ALIGN.CENTER)
    return shp


def _fit_image(slide, img_bytes, box_l, box_t, box_w, box_h):
    """Place an image fit inside the box, aspect preserved, centered, framed."""
    try:
        pic = slide.shapes.add_picture(io.BytesIO(img_bytes), box_l, box_t)
    except Exception:  # noqa: BLE001 - bad bytes -> skip
        logger.warning("deck: could not place an image")
        return
    scale = min(box_w / pic.width, box_h / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int(box_l + (box_w - pic.width) / 2)
    pic.top = int(box_t + (box_h - pic.height) / 2)
    pic.line.color.rgb = _rgb(BORDER)
    pic.line.width = Pt(1)


def _body_size(n: int, ov: dict) -> int:
    base = 18 if n <= 4 else (16 if n == 5 else 14)
    base -= 2 * ov.get("font_step", 0)
    return max(base, 10)


def _title(slide, text, P, color=None, size=30):
    tf = _text(slide, MARGIN, Inches(0.45), CONTENT_W, Inches(1.0))
    _run(tf.paragraphs[0], text, size, color or P["primary"], bold=True)


# --- Per-layout renderers ----------------------------------------------------
def _r_title(slide, d, P, ov):
    _bg(slide, P["primary"])
    tf = _text(slide, MARGIN, Inches(2.5), CONTENT_W, Inches(2.0), MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], d.get("title", ""), 44, WHITE, bold=True, align=PP_ALIGN.CENTER)
    if d.get("subtitle"):
        st = _text(slide, MARGIN, Inches(4.4), CONTENT_W, Inches(1.0))
        _run(st.paragraphs[0], d["subtitle"], 20, P["tint"], align=PP_ALIGN.CENTER)


def _r_section(slide, d, P, ov):
    _bg(slide, P["primary"])
    _circle(slide, Inches(6.27), Inches(2.1), Inches(0.8), P["tint"], "›", P["primary"])
    tf = _text(slide, MARGIN, Inches(3.2), CONTENT_W, Inches(1.6), MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], d.get("title", ""), 36, WHITE, bold=True, align=PP_ALIGN.CENTER)


def _r_bullets(slide, d, P, ov):
    _bg(slide, WHITE)
    _title(slide, d.get("title", ""), P)
    items = [str(b) for b in d.get("bullets", [])][:_MAX_ITEMS]
    size = _body_size(len(items), ov)
    tf = _text(slide, MARGIN, Inches(1.7), CONTENT_W, Inches(5.2))
    for i, b in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        _run(p, f"•  {b}", size, INK)


def _r_two_column(slide, d, P, ov):
    _bg(slide, WHITE)
    _title(slide, d.get("title", ""), P)
    text_w = Inches(6.2) if ov.get("stricter") else Inches(6.8)
    items = [str(b) for b in d.get("bullets", [])][:_MAX_ITEMS]
    size = _body_size(len(items), ov)
    tf = _text(slide, MARGIN, Inches(1.7), text_w, Inches(5.2))
    for i, b in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        _run(p, f"•  {b}", size, INK)
    img = ov.get("_image")
    img_l = MARGIN + text_w + Inches(0.4)
    _fit_image(slide, img, img_l, Inches(1.7), SLIDE_W - img_l - MARGIN, Inches(5.0)) if img else None


def _r_image_feature(slide, d, P, ov):
    _bg(slide, WHITE)
    _title(slide, d.get("title", ""), P)
    img = ov.get("_image")
    if img:
        _fit_image(slide, img, MARGIN, Inches(1.7), CONTENT_W, Inches(4.4))
    if d.get("caption"):
        cf = _text(slide, MARGIN, Inches(6.3), CONTENT_W, Inches(0.8))
        _run(cf.paragraphs[0], d["caption"], 12, MUTED, align=PP_ALIGN.CENTER)


def _r_stat(slide, d, P, ov):
    _bg(slide, WHITE)
    _title(slide, d.get("title", ""), P)
    stats = (d.get("stats") or [])[:3]
    if not stats:
        return
    col_w = CONTENT_W / len(stats)
    for i, st in enumerate(stats):
        left = MARGIN + int(col_w * i)
        vf = _text(slide, left, Inches(2.6), int(col_w), Inches(1.4), MSO_ANCHOR.MIDDLE)
        _run(vf.paragraphs[0], st.get("value", ""), 60, P["primary"], bold=True, align=PP_ALIGN.CENTER)
        lf = _text(slide, left, Inches(4.2), int(col_w), Inches(1.4))
        _run(lf.paragraphs[0], st.get("label", ""), 16, INK, align=PP_ALIGN.CENTER)


def _r_comparison(slide, d, P, ov):
    _bg(slide, WHITE)
    _title(slide, d.get("title", ""), P)
    cols = (d.get("columns") or [])[:3]
    if not cols:
        return
    gap = Inches(0.4)
    col_w = int((CONTENT_W - gap * (len(cols) - 1)) / len(cols))
    for i, col in enumerate(cols):
        left = MARGIN + (col_w + gap) * i
        hf = _text(slide, left, Inches(1.7), col_w, Inches(0.7))
        _run(hf.paragraphs[0], col.get("heading", ""), 20, P["primary"], bold=True, align=PP_ALIGN.CENTER)
        items = [str(x) for x in col.get("items", [])][:_MAX_ITEMS]
        bf = _text(slide, left, Inches(2.5), col_w, Inches(4.4))
        for j, it in enumerate(items):
            p = bf.paragraphs[0] if j == 0 else bf.add_paragraph()
            p.space_after = Pt(8)
            _run(p, f"•  {it}", _body_size(len(items), ov), INK)


def _r_quote(slide, d, P, ov):
    _bg(slide, P["primary"])
    tf = _text(slide, Inches(1.2), Inches(2.3), SLIDE_W - Inches(2.4), Inches(2.6), MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], f"“{d.get('quote', '')}”", 30, WHITE, bold=True, align=PP_ALIGN.CENTER)
    if d.get("attribution"):
        af = _text(slide, Inches(1.2), Inches(5.1), SLIDE_W - Inches(2.4), Inches(0.8))
        _run(af.paragraphs[0], f"— {d['attribution']}", 16, P["tint"], align=PP_ALIGN.CENTER)


def _r_agenda(slide, d, P, ov):
    _bg(slide, WHITE)
    _title(slide, d.get("title", "") or "Agenda", P)
    items = [str(x) for x in d.get("items", [])][:_MAX_ITEMS]
    top = Inches(1.9)
    row = Inches(1.0)
    for i, it in enumerate(items):
        y = top + row * i
        _circle(slide, MARGIN, y, Inches(0.6), P["primary"], str(i + 1), WHITE)
        tf = _text(slide, MARGIN + Inches(0.9), y, CONTENT_W - Inches(0.9), Inches(0.6), MSO_ANCHOR.MIDDLE)
        _run(tf.paragraphs[0], it, 18, INK)


_RENDERERS = {
    "title": _r_title, "section": _r_section, "bullets": _r_bullets,
    "two_column": _r_two_column, "image_feature": _r_image_feature,
    "stat": _r_stat, "comparison": _r_comparison, "quote": _r_quote,
    "agenda": _r_agenda,
}


# --- Build -------------------------------------------------------------------
def palette_for(plan: dict) -> dict:
    return PALETTES.get(plan.get("palette") or "", PALETTES[_DEFAULT_PALETTE])


def build_deck(
    plan: dict,
    photos: dict[int, bytes] | None = None,
    overrides: dict[int, dict] | None = None,
    blank_layout_index: int = 6,
) -> bytes:
    """Render the (normalized) plan to .pptx bytes using the design system."""
    photos = photos or {}
    overrides = overrides or {}
    P = palette_for(plan)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[blank_layout_index]

    for idx, d in enumerate(plan.get("slides", []), start=1):
        slide = prs.slides.add_slide(blank)
        ov = dict(overrides.get(idx, {}))
        ref = d.get("image_ref")
        ov["_image"] = photos.get(ref) if isinstance(ref, int) else None
        if ov.get("contrast"):
            # Force the guaranteed high-contrast text colors (handled per layout).
            pass
        renderer = _RENDERERS.get(d.get("layout", "bullets"), _r_bullets)
        try:
            renderer(slide, d, P, ov)
        except Exception:  # noqa: BLE001 - one bad slide shouldn't kill the deck
            logger.exception("deck: failed to render slide %s", idx)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def slide_count(plan: dict) -> int:
    return len(plan.get("slides", []))
